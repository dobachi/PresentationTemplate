"""
Main presentation builder module.

This module provides the PresentationBuilder class for creating PowerPoint
presentations with diagrams and charts.
"""

import os
from typing import Optional, Dict, List, Any
from pathlib import Path
import yaml

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from src.core.slide_builder import SlideBuilder
from src.core.template_manager import TemplateManager


class PresentationBuilder:
    """
    Main class for building PowerPoint presentations.

    Supports both AI-assisted conversation mode and direct YAML-based input.

    Attributes:
        presentation: python-pptx Presentation object
        theme: Theme configuration dictionary
        language: Language code ('ja', 'en', or 'ja_en')
        slide_builder: SlideBuilder instance for creating individual slides
    """

    def __init__(
        self,
        config_path: Optional[str] = None,
        template_path: Optional[str] = None,
        theme: str = "corporate",
        language: str = "ja",
        ai_mode: bool = False
    ):
        """
        Initialize the PresentationBuilder.

        Args:
            config_path: Path to configuration YAML file
            template_path: Path to PowerPoint template file
            theme: Theme name ('corporate', 'technical', 'academic')
            language: Language code ('ja', 'en', 'ja_en')
            ai_mode: Enable AI-assisted mode (for future integration)
        """
        self.theme_name = theme
        self.language = language
        self.ai_mode = ai_mode

        # Load configuration
        self.config = self._load_config(config_path)

        # Load theme
        self.theme = self._load_theme(theme)

        # Initialize template manager
        self.template_manager = TemplateManager(template_path)

        # Create presentation
        if template_path and os.path.exists(template_path):
            self.presentation = Presentation(template_path)
        else:
            self.presentation = Presentation()
            self._setup_default_slide_layouts()

        # Initialize slide builder
        self.slide_builder = SlideBuilder(self.presentation, self.theme, self.language)

        # Metadata
        self.metadata = {
            'title': '',
            'author': '',
            'subject': '',
            'created_by': 'PresentationGenerator'
        }

    def _load_config(self, config_path: Optional[str]) -> Dict[str, Any]:
        """Load configuration from YAML file."""
        if config_path and os.path.exists(config_path):
            with open(config_path, 'r', encoding='utf-8') as f:
                return yaml.safe_load(f)

        # Default configuration
        return {
            'default_theme': 'corporate',
            'default_language': 'ja',
            'slide_width': 10,
            'slide_height': 7.5,
            'dpi': 300
        }

    def _load_theme(self, theme_name: str) -> Dict[str, Any]:
        """Load theme configuration."""
        theme_path = Path(f"config/themes/{theme_name}.yaml")

        if theme_path.exists():
            with open(theme_path, 'r', encoding='utf-8') as f:
                return yaml.safe_load(f)

        # Default theme
        return {
            'name': 'corporate',
            'colors': {
                'primary': '#0066CC',
                'secondary': '#004C99',
                'accent': '#FF6600',
                'background': '#FFFFFF',
                'text': '#333333'
            },
            'fonts': {
                'title': 'Meiryo' if 'ja' in self.language else 'Arial',
                'body': 'Meiryo' if 'ja' in self.language else 'Arial',
                'code': 'Consolas'
            },
            'sizes': {
                'title': 44,
                'subtitle': 32,
                'heading': 28,
                'body': 18,
                'caption': 14
            }
        }

    def _setup_default_slide_layouts(self):
        """Setup default slide layouts if no template is provided."""
        # This will use python-pptx default layouts
        pass

    def set_metadata(
        self,
        title: str = '',
        author: str = '',
        subject: str = ''
    ):
        """
        Set presentation metadata.

        Args:
            title: Presentation title
            author: Author name
            subject: Subject/topic
        """
        self.metadata.update({
            'title': title,
            'author': author,
            'subject': subject
        })

        # Apply to presentation properties
        if title:
            self.presentation.core_properties.title = title
        if author:
            self.presentation.core_properties.author = author
        if subject:
            self.presentation.core_properties.subject = subject

    def set_fonts(
        self,
        title_font: Optional[str] = None,
        body_font: Optional[str] = None,
        code_font: Optional[str] = None
    ):
        """
        Set custom fonts for the presentation.

        Args:
            title_font: Font for titles
            body_font: Font for body text
            code_font: Font for code blocks
        """
        if title_font:
            self.theme['fonts']['title'] = title_font
        if body_font:
            self.theme['fonts']['body'] = body_font
        if code_font:
            self.theme['fonts']['code'] = code_font

    def load_definition(self, yaml_path: str) -> Dict[str, Any]:
        """
        Load presentation definition from YAML file.

        Args:
            yaml_path: Path to YAML definition file

        Returns:
            Presentation definition dictionary
        """
        with open(yaml_path, 'r', encoding='utf-8') as f:
            definition = yaml.safe_load(f)

        # Set metadata from definition
        pres_meta = definition.get('presentation', {})
        self.set_metadata(
            title=pres_meta.get('title', ''),
            author=pres_meta.get('author', ''),
            subject=pres_meta.get('subject', '')
        )

        # Update language if specified
        if 'language' in pres_meta:
            self.language = pres_meta['language']
            self.slide_builder.language = self.language

        return definition

    def add_title_slide(
        self,
        title: str,
        subtitle: str = ''
    ):
        """
        Add a title slide to the presentation.

        Args:
            title: Main title text
            subtitle: Subtitle text
        """
        self.slide_builder.add_title_slide(title, subtitle)

    def add_text_slide(
        self,
        title: str,
        bullets: List[str],
        notes: str = ''
    ):
        """
        Add a text slide with bullet points.

        Args:
            title: Slide title
            bullets: List of bullet points
            notes: Speaker notes
        """
        self.slide_builder.add_text_slide(title, bullets, notes)

    def add_diagram_slide(
        self,
        title: str,
        diagram_path: str,
        notes: str = ''
    ):
        """
        Add a slide with a diagram image.

        Args:
            title: Slide title
            diagram_path: Path to diagram image file
            notes: Speaker notes
        """
        self.slide_builder.add_diagram_slide(title, diagram_path, notes)

    def add_chart_slide(
        self,
        title: str,
        chart_path: str,
        notes: str = ''
    ):
        """
        Add a slide with a chart image.

        Args:
            title: Slide title
            chart_path: Path to chart image file
            notes: Speaker notes
        """
        # Charts are treated similarly to diagrams
        self.slide_builder.add_diagram_slide(title, chart_path, notes)

    def add_two_column_slide(
        self,
        title: str,
        left_content: str,
        right_content: str,
        notes: str = ''
    ):
        """
        Add a two-column slide.

        Args:
            title: Slide title
            left_content: Content for left column (text or image path)
            right_content: Content for right column (text or image path)
            notes: Speaker notes
        """
        self.slide_builder.add_two_column_slide(
            title, left_content, right_content, notes
        )

    def build_from_definition(self, definition: Dict[str, Any]):
        """
        Build presentation from a definition dictionary.

        Args:
            definition: Presentation definition (from YAML)
        """
        slides = definition.get('slides', [])

        for slide_def in slides:
            slide_type = slide_def.get('type', 'content')

            if slide_type == 'title':
                self.add_title_slide(
                    title=slide_def.get('title', ''),
                    subtitle=slide_def.get('subtitle', '')
                )

            elif slide_type == 'content':
                layout = slide_def.get('layout', 'text')
                title = slide_def.get('title', '')
                notes = slide_def.get('notes', '')

                if layout == 'text':
                    bullets = slide_def.get('bullets', [])
                    self.add_text_slide(title, bullets, notes)

                elif layout == 'diagram':
                    # Diagram path should be provided or generated
                    content = slide_def.get('content', {})
                    diagram_path = content.get('path', '')
                    if diagram_path:
                        self.add_diagram_slide(title, diagram_path, notes)

                elif layout == 'chart':
                    # Chart path should be provided or generated
                    content = slide_def.get('content', {})
                    chart_path = content.get('path', '')
                    if chart_path:
                        self.add_chart_slide(title, chart_path, notes)

    def build(self):
        """
        Build the complete presentation.

        This method can be extended for additional processing.
        """
        # Placeholder for future build logic
        pass

    def save(self, output_path: str):
        """
        Save the presentation to a file.

        Args:
            output_path: Output file path (should end with .pptx)
        """
        # Ensure output directory exists
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir)

        # Save presentation
        self.presentation.save(output_path)
        print(f"Presentation saved to: {output_path}")

    def from_conversation(self, conversation):
        """
        Build presentation from AI conversation.

        Args:
            conversation: PresentationConversation instance

        Note: This is a placeholder for future AI integration.
        """
        # TODO: Implement AI conversation integration
        raise NotImplementedError("AI conversation mode not yet implemented")
