"""
Slide builder module for creating individual slides.

This module handles the creation of different types of slides with
proper layout and styling.
"""

import os
from typing import List, Dict, Any, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class SlideBuilder:
    """
    Builder class for creating individual presentation slides.

    Attributes:
        presentation: python-pptx Presentation object
        theme: Theme configuration dictionary
        language: Language code for font selection
    """

    def __init__(
        self,
        presentation: Presentation,
        theme: Dict[str, Any],
        language: str = "ja"
    ):
        """
        Initialize the SlideBuilder.

        Args:
            presentation: python-pptx Presentation object
            theme: Theme configuration dictionary
            language: Language code ('ja', 'en', 'ja_en')
        """
        self.presentation = presentation
        self.theme = theme
        self.language = language

    def _get_blank_slide_layout(self):
        """Get blank slide layout."""
        # Layout 6 is typically blank
        return self.presentation.slide_layouts[6]

    def _get_title_slide_layout(self):
        """Get title slide layout."""
        # Layout 0 is typically title slide
        return self.presentation.slide_layouts[0]

    def _get_title_content_layout(self):
        """Get title and content layout."""
        # Layout 1 is typically title and content
        return self.presentation.slide_layouts[1]

    def _apply_font(self, text_frame, font_name: str, font_size: int, bold: bool = False):
        """Apply font settings to a text frame."""
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.bold = bold

    def _parse_color(self, color_str: str) -> RGBColor:
        """
        Parse color string to RGBColor.

        Args:
            color_str: Color in hex format (#RRGGBB)

        Returns:
            RGBColor object
        """
        color_str = color_str.lstrip('#')
        r, g, b = tuple(int(color_str[i:i+2], 16) for i in (0, 2, 4))
        return RGBColor(r, g, b)

    def add_title_slide(self, title: str, subtitle: str = ''):
        """
        Add a title slide.

        Args:
            title: Main title text
            subtitle: Subtitle text
        """
        slide_layout = self._get_title_slide_layout()
        slide = self.presentation.slides.add_slide(slide_layout)

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = title
            title_frame = slide.shapes.title.text_frame
            self._apply_font(
                title_frame,
                self.theme['fonts']['title'],
                self.theme['sizes']['title'],
                bold=True
            )

        # Set subtitle if present
        if len(slide.placeholders) > 1 and subtitle:
            slide.placeholders[1].text = subtitle
            subtitle_frame = slide.placeholders[1].text_frame
            self._apply_font(
                subtitle_frame,
                self.theme['fonts']['body'],
                self.theme['sizes']['subtitle']
            )

    def add_text_slide(
        self,
        title: str,
        bullets: List[str],
        notes: str = ''
    ):
        """
        Add a text slide with bullet points.

        Args:
            title: Slide title
            bullets: List of bullet point strings
            notes: Speaker notes
        """
        slide_layout = self._get_title_content_layout()
        slide = self.presentation.slides.add_slide(slide_layout)

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = title
            title_frame = slide.shapes.title.text_frame
            self._apply_font(
                title_frame,
                self.theme['fonts']['title'],
                self.theme['sizes']['heading'],
                bold=True
            )

        # Add content
        if len(slide.placeholders) > 1:
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()

            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()

                p.text = bullet
                p.level = 0
                p.font.name = self.theme['fonts']['body']
                p.font.size = Pt(self.theme['sizes']['body'])

        # Add notes if provided
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_diagram_slide(
        self,
        title: str,
        diagram_path: str,
        notes: str = ''
    ):
        """
        Add a slide with a diagram image.

        Args:
            title: Slide title
            diagram_path: Path to diagram image file
            notes: Speaker notes
        """
        slide_layout = self._get_blank_slide_layout()
        slide = self.presentation.slides.add_slide(slide_layout)

        # Add title text box
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.8)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title

        # Style title
        p = title_frame.paragraphs[0]
        p.font.name = self.theme['fonts']['title']
        p.font.size = Pt(self.theme['sizes']['heading'])
        p.font.bold = True
        p.font.color.rgb = self._parse_color(self.theme['colors']['primary'])

        # Add diagram image
        if os.path.exists(diagram_path):
            # Calculate position to center the image
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(5)

            slide.shapes.add_picture(
                diagram_path,
                left, top,
                width=width,
                height=height
            )

        # Add notes if provided
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_two_column_slide(
        self,
        title: str,
        left_content: str,
        right_content: str,
        notes: str = ''
    ):
        """
        Add a two-column slide.

        Args:
            title: Slide title
            left_content: Content for left column (text or image path)
            right_content: Content for right column (text or image path)
            notes: Speaker notes
        """
        slide_layout = self._get_blank_slide_layout()
        slide = self.presentation.slides.add_slide(slide_layout)

        # Add title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.8)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title

        p = title_frame.paragraphs[0]
        p.font.name = self.theme['fonts']['title']
        p.font.size = Pt(self.theme['sizes']['heading'])
        p.font.bold = True

        # Left column
        left_col_left = Inches(0.5)
        left_col_top = Inches(1.5)
        left_col_width = Inches(4.5)
        left_col_height = Inches(5)

        if os.path.exists(left_content):
            # It's an image
            slide.shapes.add_picture(
                left_content,
                left_col_left, left_col_top,
                width=left_col_width
            )
        else:
            # It's text
            left_box = slide.shapes.add_textbox(
                left_col_left, left_col_top,
                left_col_width, left_col_height
            )
            left_frame = left_box.text_frame
            left_frame.text = left_content
            self._apply_font(
                left_frame,
                self.theme['fonts']['body'],
                self.theme['sizes']['body']
            )

        # Right column
        right_col_left = Inches(5.5)
        right_col_top = Inches(1.5)
        right_col_width = Inches(4.5)
        right_col_height = Inches(5)

        if os.path.exists(right_content):
            # It's an image
            slide.shapes.add_picture(
                right_content,
                right_col_left, right_col_top,
                width=right_col_width
            )
        else:
            # It's text
            right_box = slide.shapes.add_textbox(
                right_col_left, right_col_top,
                right_col_width, right_col_height
            )
            right_frame = right_box.text_frame
            right_frame.text = right_content
            self._apply_font(
                right_frame,
                self.theme['fonts']['body'],
                self.theme['sizes']['body']
            )

        # Add notes if provided
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes
